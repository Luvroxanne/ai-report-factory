from __future__ import annotations

from pathlib import Path
from typing import Any

from app.adapters.presenton_client import PresentonClient


class PPTService:
    def __init__(self, presenton: PresentonClient):
        self.presenton = presenton

    def generate(
        self,
        *,
        task_id: str,
        text: str,
        plan: dict[str, Any],
        output_dir: Path,
        logger: Any | None = None,
    ) -> Path:
        output_path = output_dir / "report.pptx"
        slides_markdown = self._slides_to_markdown(plan)
        instructions = (
            "生成正式商业 PPT，16:9 宽屏，包含封面、目录、章节、内容与总结页；"
            "使用统一主题色、字体、留白、图标感和卡片化布局。"
        )

        try:
            generated = self.presenton.generate_presentation(
                content=text,
                slides_markdown=slides_markdown,
                instructions=instructions,
                output_path=output_path,
            )
            if generated:
                if logger:
                    logger.info("Presenton 生成 PPT", path=generated)
                return generated
        except Exception as exc:
            if logger:
                logger.warning("Presenton 不可用，切换本地 PPT 兜底", reason=exc)

        path = self._generate_local_ppt(plan, output_path)
        if logger:
            logger.info("本地 PPT 兜底完成", path=path)
        return path

    def _slides_to_markdown(self, plan: dict[str, Any]) -> list[str]:
        pages: list[str] = []
        for slide in plan.get("slides", []):
            lines = [f"# {slide.get('title', '')}"]
            if slide.get("chapter"):
                lines.append(f"章节：{slide.get('chapter')}")
            for bullet in slide.get("bullets", []):
                lines.append(f"- {bullet}")
            if slide.get("speaker_note"):
                lines.append(f"讲稿：{slide.get('speaker_note')}")
            pages.append("\n".join(lines))
        return pages

    def _generate_local_ppt(self, plan: dict[str, Any], output_path: Path) -> Path:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slides = plan.get("slides", []) or []
        if not slides:
            slides = [{"layout": "cover", "title": plan.get("title") or "AI报告", "bullets": []}]

        theme = self._theme(plan.get("style", "official-tech"))
        for index, slide_data in enumerate(slides, start=1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            layout = slide_data.get("layout") or "content"
            self._paint_background(slide, prs, theme, index)
            if layout == "cover" or index == 1:
                self._render_cover(slide, plan, slide_data, theme)
            elif layout == "agenda":
                self._render_agenda(slide, slide_data, theme, index, len(slides))
            elif layout == "section":
                self._render_section(slide, slide_data, theme, index, len(slides))
            elif layout == "summary" or index == len(slides):
                self._render_summary(slide, slide_data, theme, index, len(slides))
            else:
                self._render_content(slide, slide_data, theme, index, len(slides))
            self._write_notes(slide, slide_data.get("speaker_note", ""))

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        return output_path

    def _theme(self, style: str) -> dict[str, tuple[int, int, int]]:
        themes = {
            "training": {
                "bg": (249, 247, 242), "panel": (255, 255, 255), "text": (32, 44, 60),
                "muted": (95, 111, 132), "accent": (205, 126, 43), "accent2": (40, 118, 152), "line": (228, 221, 210),
            },
            "roadshow": {
                "bg": (13, 12, 26), "panel": (29, 26, 48), "text": (252, 247, 232),
                "muted": (188, 178, 158), "accent": (242, 194, 88), "accent2": (232, 85, 74), "line": (81, 69, 105),
            },
        }
        return themes.get(style, {
            "bg": (5, 17, 34), "panel": (13, 34, 64), "text": (244, 248, 255),
            "muted": (155, 179, 210), "accent": (47, 214, 188), "accent2": (47, 140, 255), "line": (62, 122, 196),
        })

    def _rgb(self, value: tuple[int, int, int]):
        from pptx.dml.color import RGBColor
        return RGBColor(*value)

    def _paint_background(self, slide: Any, prs: Any, theme: dict[str, tuple[int, int, int]], index: int) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(theme["bg"])
        bg.line.fill.background()

        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
        band.fill.solid()
        band.fill.fore_color.rgb = self._rgb(theme["accent2"])
        band.line.fill.background()

        halo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.0), Inches(-0.8), Inches(4.9), Inches(4.9))
        halo.fill.solid()
        halo.fill.fore_color.rgb = self._rgb(theme["accent2"])
        halo.fill.transparency = 35
        halo.line.fill.background()

        halo2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.85), Inches(5.35), Inches(1.45), Inches(1.45))
        halo2.fill.solid()
        halo2.fill.fore_color.rgb = self._rgb(theme["accent"])
        halo2.fill.transparency = 22
        halo2.line.fill.background()

    def _add_text(self, slide: Any, text: str, left: float, top: float, width: float, height: float, size: int, theme: dict[str, tuple[int, int, int]], *, bold: bool = False, color: str = "text", align: Any | None = None) -> Any:
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = self._rgb(theme[color])
        if align is not None:
            p.alignment = align
        return box

    def _add_card(self, slide: Any, left: float, top: float, width: float, height: float, theme: dict[str, tuple[int, int, int]], transparency: int = 4) -> Any:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = self._rgb(theme["panel"])
        card.fill.transparency = transparency
        card.line.color.rgb = self._rgb(theme["line"])
        card.line.transparency = 25
        return card

    def _render_cover(self, slide: Any, plan: dict[str, Any], slide_data: dict[str, Any], theme: dict[str, tuple[int, int, int]]) -> None:
        from pptx.enum.text import PP_ALIGN

        title = plan.get("title") or slide_data.get("title") or "AI报告"
        subtitle = plan.get("subtitle") or "PPT · 解说稿 · 语音 · 1080P 视频一体生成"
        summary = plan.get("summary") or "支持 AI Provider 增强，也支持无 Token 本地兜底完整交付。"

        self._add_text(slide, "AI REPORT FACTORY", 0.82, 0.72, 5.4, 0.32, 10, theme, bold=True, color="accent")
        self._add_text(slide, title, 0.82, 1.36, 8.9, 1.42, 38, theme, bold=True)
        self._add_text(slide, subtitle, 0.88, 2.84, 8.0, 0.52, 15, theme, color="accent")
        self._add_card(slide, 0.88, 3.76, 7.7, 1.02, theme, transparency=8)
        self._add_text(slide, summary[:120] + ("…" if len(summary) > 120 else ""), 1.12, 3.96, 7.25, 0.58, 12, theme, color="muted")

        tags = ["Presenton 优先", "DOCX 讲稿", "Windows TTS", "1080P H.264", "可配置"]
        for i, tag in enumerate(tags):
            self._add_card(slide, 0.88 + i * 1.42, 5.45, 1.22, 0.38, theme, transparency=0)
            self._add_text(slide, tag, 0.9 + i * 1.42, 5.53, 1.18, 0.14, 7, theme, bold=True, color="accent", align=PP_ALIGN.CENTER)
        self._add_text(slide, "可交付 · 无 Token 也可运行", 8.72, 6.08, 3.82, 0.38, 13, theme, bold=True, color="accent", align=PP_ALIGN.RIGHT)

    def _render_agenda(self, slide: Any, slide_data: dict[str, Any], theme: dict[str, tuple[int, int, int]], index: int, total: int) -> None:
        self._header(slide, slide_data.get("title", "目录"), theme, index, total)
        bullets = slide_data.get("bullets", [])[:5]
        for i, item in enumerate(bullets, start=1):
            y = 1.65 + (i - 1) * 0.9
            self._add_card(slide, 1.0, y, 10.9, 0.62, theme, transparency=6)
            self._add_text(slide, f"{i:02d}", 1.28, y + 0.14, 0.6, 0.2, 13, theme, bold=True, color="accent")
            self._add_text(slide, str(item), 2.05, y + 0.12, 8.7, 0.25, 16, theme, bold=True)

    def _render_section(self, slide: Any, slide_data: dict[str, Any], theme: dict[str, tuple[int, int, int]], index: int, total: int) -> None:
        chapter = slide_data.get("chapter") or f"第 {index} 部分"
        title = slide_data.get("title", "章节")
        self._add_text(slide, chapter.upper(), 0.92, 1.1, 4.4, 0.36, 14, theme, bold=True, color="accent")
        self._add_text(slide, title, 0.92, 1.72, 8.6, 1.02, 36, theme, bold=True)
        self._add_card(slide, 0.96, 3.18, 9.2, 1.35, theme, transparency=8)
        bullets = slide_data.get("bullets", [])[:3]
        text = " / ".join(str(item) for item in bullets) or "聚焦关键目标、核心内容与落地路径"
        self._add_text(slide, text, 1.22, 3.56, 8.65, 0.42, 15, theme, color="muted")
        self._footer(slide, theme, index, total)

    def _render_content(self, slide: Any, slide_data: dict[str, Any], theme: dict[str, tuple[int, int, int]], index: int, total: int) -> None:
        self._header(slide, slide_data.get("title", "内容页"), theme, index, total)
        bullets = [str(item) for item in slide_data.get("bullets", [])[:6]] or ["明确背景与目标", "梳理关键路径", "形成交付闭环"]
        for i, bullet in enumerate(bullets):
            col = i % 2
            row = i // 2
            left = 0.82 + col * 6.05
            top = 1.72 + row * 1.38
            self._bullet_card(slide, left, top, 5.55, 1.06, i + 1, bullet, theme)
        self._footer(slide, theme, index, total)

    def _render_summary(self, slide: Any, slide_data: dict[str, Any], theme: dict[str, tuple[int, int, int]], index: int, total: int) -> None:
        self._header(slide, slide_data.get("title", "总结与下一步"), theme, index, total)
        bullets = [str(item) for item in slide_data.get("bullets", [])[:4]] or ["目标清晰", "路径明确", "风险可控", "持续优化"]
        for i, bullet in enumerate(bullets):
            self._add_card(slide, 0.98, 1.58 + i * 1.0, 10.85, 0.66, theme, transparency=6)
            self._add_text(slide, "✓", 1.24, 1.72 + i * 1.0, 0.35, 0.2, 15, theme, bold=True, color="accent")
            self._add_text(slide, bullet, 1.78, 1.72 + i * 1.0, 8.8, 0.26, 16, theme, bold=True)
        self._add_text(slide, "谢谢观看", 8.72, 6.0, 3.4, 0.48, 24, theme, bold=True, color="accent")

    def _header(self, slide: Any, title: str, theme: dict[str, tuple[int, int, int]], index: int, total: int) -> None:
        from pptx.enum.text import PP_ALIGN
        self._add_text(slide, title, 0.72, 0.52, 10.1, 0.72, 26, theme, bold=True)
        self._add_text(slide, f"{index:02d} / {total:02d}", 11.1, 0.6, 1.4, 0.22, 10, theme, bold=True, color="accent", align=PP_ALIGN.RIGHT)

    def _footer(self, slide: Any, theme: dict[str, tuple[int, int, int]], index: int, total: int) -> None:
        from pptx.enum.text import PP_ALIGN
        self._add_text(slide, "AI报告工厂 · PPT / 讲稿 / 语音 / 视频 统一归档", 0.72, 6.95, 11.8, 0.2, 8, theme, color="muted", align=PP_ALIGN.RIGHT)

    def _bullet_card(self, slide: Any, left: float, top: float, width: float, height: float, number: int, text: str, theme: dict[str, tuple[int, int, int]]) -> None:
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        self._add_card(slide, left, top, width, height, theme, transparency=5)
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.22), Inches(top + 0.24), Inches(0.48), Inches(0.48))
        num.fill.solid()
        num.fill.fore_color.rgb = self._rgb(theme["accent"])
        num.line.fill.background()
        p = num.text_frame.paragraphs[0]
        p.text = str(number)
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Microsoft YaHei"
        from pptx.util import Pt
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self._rgb(theme["bg"])
        self._add_text(slide, text[:86] + ("…" if len(text) > 86 else ""), left + 0.88, top + 0.22, width - 1.15, height - 0.28, 13, theme)

    def _write_notes(self, slide: Any, note: str) -> None:
        if not note:
            return
        try:
            notes = slide.notes_slide.notes_text_frame
            notes.text = note
        except Exception:
            return
